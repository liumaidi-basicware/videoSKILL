#!/usr/bin/env python3
"""Universal document text extractor — dispatch by file extension.

Supported: .pptx .pdf .docx .doc .rtf .txt .md .xlsx .csv
Strategy per type (all LOCAL, no API):
  pptx -> python-pptx
  pdf  -> PyMuPDF (fitz)
  docx -> python-docx
  doc/.rtf -> macOS textutil (fallback) ; antiword if present
  xlsx -> openpyxl
  txt/md/csv -> plain read

Returns a list of {"page": int, "text": str} blocks so callers can keep source
locality. Raises SystemExit with a clear message on unsupported/missing deps.
"""
import os
import subprocess
import shutil


def _from_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    blocks = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs).strip()
                    if line:
                        blocks.append({"page": si, "text": line})
    return blocks


def _from_pdf(path):
    import fitz
    blocks = []
    doc = fitz.open(path)
    for pi in range(len(doc)):
        text = doc[pi].get_text().strip()
        for line in text.splitlines():
            line = line.strip()
            if line:
                blocks.append({"page": pi + 1, "text": line})
    doc.close()
    return blocks


def _from_docx(path):
    from docx import Document
    d = Document(path)
    blocks = [{"page": 1, "text": p.text.strip()} for p in d.paragraphs if p.text.strip()]
    # tables too
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                blocks.append({"page": 1, "text": " | ".join(cells)})
    return blocks


def _from_textutil(path):
    """macOS: convert legacy .doc / .rtf to txt via textutil."""
    if not shutil.which("textutil"):
        raise SystemExit("legacy .doc/.rtf needs macOS 'textutil' (not found). "
                         "Convert to .docx/.pdf first, or install antiword.")
    out = subprocess.run(["textutil", "-convert", "txt", "-stdout", path],
                         stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    if out.returncode != 0:
        raise SystemExit("textutil failed: " + out.stderr.decode("utf-8", "replace")[:300])
    text = out.stdout.decode("utf-8", "replace")
    return [{"page": 1, "text": ln.strip()} for ln in text.splitlines() if ln.strip()]


def _from_xlsx(path):
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    blocks = []
    for si, ws in enumerate(wb.worksheets, 1):
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                blocks.append({"page": si, "text": " | ".join(cells)})
    return blocks


def _from_plain(path):
    with open(path, encoding="utf-8", errors="replace") as f:
        return [{"page": 1, "text": ln.strip()} for ln in f if ln.strip()]


_DISPATCH = {
    ".pptx": _from_pptx,
    ".pdf": _from_pdf,
    ".docx": _from_docx,
    ".doc": _from_textutil,
    ".rtf": _from_textutil,
    ".xlsx": _from_xlsx,
    ".txt": _from_plain,
    ".md": _from_plain,
    ".csv": _from_plain,
}

SUPPORTED = sorted(_DISPATCH.keys())


def extract(path):
    """Return list of {page, text} blocks. Raises SystemExit on unsupported type."""
    if not os.path.isfile(path):
        raise SystemExit("file not found: %s" % path)
    ext = os.path.splitext(path)[1].lower()
    fn = _DISPATCH.get(ext)
    if not fn:
        raise SystemExit("unsupported file type '%s'. Supported: %s"
                         % (ext, ", ".join(SUPPORTED)))
    return fn(path)


def count_images(path):
    """Best-effort embedded image count (pptx/pdf/docx)."""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            return sum(1 for s in prs.slides for sh in s.shapes
                       if getattr(sh, "shape_type", None) == 13)
        if ext == ".pdf":
            import fitz
            doc = fitz.open(path)
            n = sum(len(doc[i].get_images()) for i in range(len(doc)))
            doc.close()
            return n
        if ext == ".docx":
            import zipfile
            with zipfile.ZipFile(path) as z:
                return sum(1 for n in z.namelist() if n.startswith("word/media/"))
    except Exception:
        return 0
    return 0


if __name__ == "__main__":
    import sys, json
    if len(sys.argv) < 2:
        print("usage: doc_extract.py <file>  (supported: %s)" % ", ".join(SUPPORTED))
        raise SystemExit(2)
    blocks = extract(sys.argv[1])
    print(json.dumps({"blocks": len(blocks), "images": count_images(sys.argv[1]),
                      "sample": blocks[:10]}, ensure_ascii=False, indent=2))
